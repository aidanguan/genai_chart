#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""验证PPTX文件"""

import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation

pptx_path = r"c:\AI\genai_chart-1\backend\temp\exports\infographic_1764236291567.pptx"

prs = Presentation(pptx_path)
slide = prs.slides[0]

print("="*60)
print(f"幻灯片尺寸: {prs.slide_width/914400:.2f}\" x {prs.slide_height/914400:.2f}\"")
print(f"比例: {(prs.slide_width/prs.slide_height):.2f}:1")
print("="*60)

text_shapes = [s for s in slide.shapes if hasattr(s, 'text_frame') and s.text]
print(f"\n文本框总数: {len(text_shapes)}\n")

for i, shape in enumerate(text_shapes[:5]):
    para = shape.text_frame.paragraphs[0]
    font = para.font
    
    print(f"{i+1}. 文本: \"{shape.text[:30]}\"")
    print(f"   字体名称: {font.name}")
    print(f"   字体大小: {font.size/12700:.1f} pt") if font.size else print("   字体大小: N/A")
    print(f"   是否粗体: {font.bold}")
    print(f"   颜色类型: {font.color.type if font.color else 'N/A'}")
    
    try:
        if font.color and hasattr(font.color, 'rgb'):
            rgb = font.color.rgb
            print(f"   RGB颜色: ({rgb[0]}, {rgb[1]}, {rgb[2]})")
    except:
        print("   RGB颜色: N/A")
    
    print()

print("="*60)
print("验证完成!")
