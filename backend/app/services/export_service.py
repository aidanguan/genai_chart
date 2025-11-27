"""
导出服务 - 支持SVG、PNG、PDF、PPTX格式导出
"""
import os
import base64
from io import BytesIO
from typing import Optional
from pathlib import Path

# 导出格式常量
EXPORT_FORMAT_SVG = "svg"
EXPORT_FORMAT_PNG = "png"
EXPORT_FORMAT_PDF = "pdf"
EXPORT_FORMAT_PPTX = "pptx"

SUPPORTED_FORMATS = [EXPORT_FORMAT_SVG, EXPORT_FORMAT_PNG, EXPORT_FORMAT_PDF, EXPORT_FORMAT_PPTX]


class ExportService:
    """导出服务类"""
    
    def __init__(self):
        self.temp_dir = Path("temp/exports")
        self.temp_dir.mkdir(parents=True, exist_ok=True)
    
    def export_svg(self, svg_content: str, filename: Optional[str] = None) -> dict:
        """
        导出SVG格式
        
        Args:
            svg_content: SVG内容字符串
            filename: 可选的文件名
            
        Returns:
            包含文件信息的字典
        """
        if not filename:
            filename = "infographic.svg"
        
        filepath = self.temp_dir / filename
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(svg_content)
        
        return {
            "format": EXPORT_FORMAT_SVG,
            "filename": filename,
            "filepath": str(filepath),
            "size": os.path.getsize(filepath)
        }
    
    def export_png(self, svg_content: str, filename: Optional[str] = None, 
                   width: int = 800, height: int = 600, scale: int = 2) -> dict:
        """
        导出PNG格式
        需要安装: pip install cairosvg
        
        Args:
            svg_content: SVG内容字符串
            filename: 可选的文件名
            width: 输出宽度
            height: 输出高度
            scale: 缩放比例(用于高清输出)
            
        Returns:
            包含文件信息的字典
        """
        try:
            import cairosvg
        except ImportError:
            raise ImportError("PNG导出需要安装cairosvg: pip install cairosvg")
        
        if not filename:
            filename = "infographic.png"
        
        filepath = self.temp_dir / filename
        
        # 转换SVG到PNG
        cairosvg.svg2png(
            bytestring=svg_content.encode('utf-8'),
            write_to=str(filepath),
            output_width=width * scale,
            output_height=height * scale
        )
        
        return {
            "format": EXPORT_FORMAT_PNG,
            "filename": filename,
            "filepath": str(filepath),
            "size": os.path.getsize(filepath),
            "width": width * scale,
            "height": height * scale
        }
    
    def export_pdf(self, svg_content: str, filename: Optional[str] = None) -> dict:
        """
        导出PDF格式
        需要安装: pip install cairosvg
        
        Args:
            svg_content: SVG内容字符串
            filename: 可选的文件名
            
        Returns:
            包含文件信息的字典
        """
        try:
            import cairosvg
        except ImportError:
            raise ImportError("PDF导出需要安装cairosvg: pip install cairosvg")
        
        if not filename:
            filename = "infographic.pdf"
        
        filepath = self.temp_dir / filename
        
        # 转换SVG到PDF
        cairosvg.svg2pdf(
            bytestring=svg_content.encode('utf-8'),
            write_to=str(filepath)
        )
        
        return {
            "format": EXPORT_FORMAT_PDF,
            "filename": filename,
            "filepath": str(filepath),
            "size": os.path.getsize(filepath)
        }
    
    def export_pptx(self, svg_content: str, title: str = "信息图", 
                    filename: Optional[str] = None) -> dict:
        """
        导出PPTX格式 - 直接插入SVG矢量图，保持可编辑性和矢量特性
        需要安装: pip install python-pptx lxml
        
        Args:
            svg_content: SVG内容字符串
            title: 幻灯片标题（可选，不会显示在幻灯片上）
            filename: 可选的文件名
            
        Returns:
            包含文件信息的字典
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Emu
            from xml.etree import ElementTree as ET
            import re
            import zipfile
            import shutil
        except ImportError:
            raise ImportError("PPTX导出需要安装: pip install python-pptx lxml")
        
        if not filename:
            filename = "infographic.pptx"
        
        # 🔄 转换 SVG 为 PPT 兼容格式 (foreignObject → <text>)
        svg_content = self._convert_svg_for_ppt(svg_content)
        
        # 解析SVG获取尺寸
        try:
            svg_root = ET.fromstring(svg_content)
        except:
            svg_content_clean = re.sub(r'xmlns="[^"]*"', '', svg_content, count=1)
            svg_root = ET.fromstring(svg_content_clean)
        
        # 获取SVG尺寸
        viewbox = svg_root.get('viewBox')
        if viewbox:
            _, _, svg_width_str, svg_height_str = viewbox.split()
            svg_width = float(svg_width_str)
            svg_height = float(svg_height_str)
        else:
            svg_width = float(svg_root.get('width', '800').replace('px', '').replace('pt', ''))
            svg_height = float(svg_root.get('height', '600').replace('px', '').replace('pt', ''))
        
        # 创建PowerPoint演示文稿
        prs = Presentation()
        
        # 设置幻灯片尺寸为标准16:9
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        
        # 添加空白幻灯片
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 计算图片位置（居中）
        slide_width_inches = 9.5
        slide_height_inches = 5.125
        
        # 计算缩放比例以适应幻灯片
        svg_aspect_ratio = svg_width / svg_height
        scale_x = (slide_width_inches * 96) / svg_width
        scale_y = (slide_height_inches * 96) / svg_height
        scale = min(scale_x, scale_y)
        
        # 计算实际显示尺寸
        display_width = (svg_width * scale) / 96
        display_height = (svg_height * scale) / 96
        
        # 居中位置
        left = Inches((10 - display_width) / 2)
        top = Inches((5.625 - display_height) / 2)
        width = Inches(display_width)
        height = Inches(display_height)
        
        # 先保存临时PPTX文件
        temp_pptx = self.temp_dir / f"temp_{filename}"
        prs.save(str(temp_pptx))
        
        # 插入SVG到PPTX
        final_filepath = self.temp_dir / filename
        self._insert_svg_to_pptx(
            str(temp_pptx),
            str(final_filepath),
            svg_content,
            left, top, width, height
        )
        
        # 删除临时文件
        if temp_pptx.exists():
            temp_pptx.unlink()
        
        return {
            "format": EXPORT_FORMAT_PPTX,
            "filename": filename,
            "filepath": str(final_filepath),
            "size": os.path.getsize(final_filepath)
        }
    
    def export(self, svg_content: str, format: str, **kwargs) -> dict:
        """
        统一导出接口
        
        Args:
            svg_content: SVG内容字符串
            format: 导出格式 (svg/png/pdf/pptx)
            **kwargs: 格式特定的参数
            
        Returns:
            包含文件信息的字典
        """
        if format not in SUPPORTED_FORMATS:
            raise ValueError(f"不支持的导出格式: {format}. 支持的格式: {SUPPORTED_FORMATS}")
        
        if format == EXPORT_FORMAT_SVG:
            return self.export_svg(svg_content, **kwargs)
        elif format == EXPORT_FORMAT_PNG:
            return self.export_png(svg_content, **kwargs)
        elif format == EXPORT_FORMAT_PDF:
            return self.export_pdf(svg_content, **kwargs)
        elif format == EXPORT_FORMAT_PPTX:
            return self.export_pptx(svg_content, **kwargs)
    
    def get_base64(self, filepath: str) -> str:
        """
        将文件转换为base64编码
        
        Args:
            filepath: 文件路径
            
        Returns:
            base64编码的字符串
        """
        with open(filepath, 'rb') as f:
            return base64.b64encode(f.read()).decode('utf-8')
    
    def _convert_svg_for_ppt(self, svg_content: str) -> str:
        """
        转换 SVG 为 PowerPoint 兼容格式
        将 foreignObject + span 转换为标准 SVG <text> 元素
        
        Args:
            svg_content: 原始 SVG 内容
            
        Returns:
            转换后的 SVG 内容
        """
        from xml.etree import ElementTree as ET
        import re
        
        try:
            # 解析 SVG
            root = ET.fromstring(svg_content)
        except:
            # 如果解析失败，尝试移除命名空间
            svg_content_clean = re.sub(r'xmlns="[^"]*"', '', svg_content, count=1)
            root = ET.fromstring(svg_content_clean)
        
        # 查找所有 foreignObject 元素
        namespaces = {'svg': 'http://www.w3.org/2000/svg', 'xhtml': 'http://www.w3.org/1999/xhtml'}
        foreign_objects = []
        
        # 递归查找 foreignObject
        def find_foreign_objects(element, path=[]):
            tag = element.tag.split('}')[-1] if '}' in element.tag else element.tag
            
            if tag == 'foreignObject':
                foreign_objects.append((element, path[:]))
            
            for i, child in enumerate(element):
                find_foreign_objects(child, path + [i])
        
        find_foreign_objects(root)
        
        # 转换每个 foreignObject
        for foreign_object, path in reversed(foreign_objects):  # 从后往前处理，避免索引变化
            text_element = self._convert_foreign_object_to_text(foreign_object)
            if text_element is not None:
                # 替换元素
                parent = root
                for idx in path[:-1]:
                    parent = list(parent)[idx]
                parent_list = list(parent)
                parent_list[path[-1]] = text_element
                parent[:] = parent_list
        
        # 转回字符串
        return ET.tostring(root, encoding='unicode')
    
    def _convert_foreign_object_to_text(self, foreign_object):
        """
        将单个 foreignObject 元素转换为 <text> 元素
        
        Args:
            foreign_object: foreignObject ET.Element
            
        Returns:
            text ET.Element 或 None
        """
        from xml.etree import ElementTree as ET
        import re
        
        # 查找 span 元素
        span = None
        for child in foreign_object:
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag == 'span':
                span = child
                break
        
        if span is None:
            return None
        
        # 提取文本内容
        text_content = span.text or ''
        if not text_content.strip():
            return None
        
        # 提取位置和尺寸
        x = float(foreign_object.get('x', '0'))
        y = float(foreign_object.get('y', '0'))
        width = float(foreign_object.get('width', '0'))
        height = float(foreign_object.get('height', '0'))
        
        # 提取样式
        style = span.get('style', '')
        styles = self._parse_style(style)
        
        # 提取字体属性
        font_size = self._extract_number(styles.get('font-size', '14px'))
        color = styles.get('color', 'rgb(38, 38, 38)')
        fill_color = self._rgb_to_hex(color)
        font_weight = styles.get('font-weight', 'normal')
        font_family = styles.get('font-family', 'inherit')
        
        # 将字体映射到系统中文字体，确保中文正常显示
        font_family = self._map_to_system_font(font_family)
        
        text_align = styles.get('text-align', 'left')
        justify_content = styles.get('justify-content', 'flex-start')
        align_items = styles.get('align-items', 'flex-start')
        
        # 计算 text-anchor
        text_anchor = 'start'
        text_x = x
        if text_align == 'center' or justify_content == 'center':
            text_anchor = 'middle'
            text_x = x + width / 2
        elif text_align == 'right' or justify_content == 'flex-end':
            text_anchor = 'end'
            text_x = x + width
        
        # 计算 dominant-baseline
        dominant_baseline = 'text-before-edge'
        text_y = y
        if align_items == 'center':
            dominant_baseline = 'middle'
            text_y = y + height / 2
        elif align_items == 'flex-end':
            dominant_baseline = 'text-after-edge'
            text_y = y + height
        
        # 创建 text 元素
        text_elem = ET.Element('text', {
            'x': str(text_x),
            'y': str(text_y),
            'fill': fill_color,
            'font-size': str(font_size),
            'font-weight': 'bold' if font_weight == 'bold' or self._extract_number(font_weight) >= 700 else 'normal',
            'font-family': font_family,
            'text-anchor': text_anchor,
            'dominant-baseline': dominant_baseline,
        })
        text_elem.text = text_content
        
        return text_elem
    
    def _parse_style(self, style_str: str) -> dict:
        """解析 CSS 样式字符串"""
        styles = {}
        for item in style_str.split(';'):
            if ':' in item:
                key, value = item.split(':', 1)
                styles[key.strip()] = value.strip()
        return styles
    
    def _extract_number(self, value: str) -> float:
        """从字符串中提取数字"""
        import re
        match = re.search(r'([0-9.]+)', str(value))
        if match:
            return float(match.group(1))
        return 14.0  # 默认值
    
    def _rgb_to_hex(self, rgb: str) -> str:
        """将 RGB 颜色转换为十六进制"""
        import re
        
        # 已经是十六进制
        if rgb.startswith('#'):
            return rgb
        
        # 解析 rgb(r, g, b) 或 rgba(r, g, b, a)
        match = re.match(r'rgba?\((\d+),\s*(\d+),\s*(\d+)', rgb)
        if not match:
            return rgb
        
        r = int(match.group(1))
        g = int(match.group(2))
        b = int(match.group(3))
        
        return f'#{r:02x}{g:02x}{b:02x}'
    
    def _map_to_system_font(self, font_family: str) -> str:
        """
        将Web字体映射到系统中文字体
        确保中文字符能够正常显示
        
        Args:
            font_family: 原始字体名称
            
        Returns:
            系统字体名称
        """
        # 常见中文字体映射
        font_mapping = {
            'Alibaba PuHuiTi': 'Noto Sans CJK SC',
            'Source Han Sans': 'Noto Sans CJK SC',
            'Source Han Serif': 'Noto Serif CJK SC',
            'LXGW WenKai': 'Noto Sans CJK SC',
            '851tegakizatsu': 'Noto Sans CJK SC',
            '黑体': 'Noto Sans CJK SC',
            '宋体': 'Noto Serif CJK SC',
            '楷体': 'Noto Sans CJK SC',
            '手写体': 'Noto Sans CJK SC',
            'sans-serif': 'Noto Sans CJK SC',
            'serif': 'Noto Serif CJK SC',
            'inherit': 'Noto Sans CJK SC',
        }
        
        # 检查是否在映射表中
        for key, value in font_mapping.items():
            if key.lower() in font_family.lower():
                return value
        
        # 如果没有匹配，默认返回 Noto Sans CJK SC
        return 'Noto Sans CJK SC'
    
    def _insert_svg_to_pptx(self, input_pptx: str, output_pptx: str, 
                           svg_content: str, left, top, width, height):
        """
        将SVG文件插入到PPTX中
        通过直接操作PPTX的ZIP结构来实现
        
        Args:
            input_pptx: 输入PPTX文件路径
            output_pptx: 输出PPTX文件路径
            svg_content: SVG内容
            left, top, width, height: 位置和尺寸
        """
        import zipfile
        import shutil
        from pathlib import Path
        import tempfile
        from lxml import etree
        
        # 创建临时目录
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            
            # 解压PPTX
            with zipfile.ZipFile(input_pptx, 'r') as zip_ref:
                zip_ref.extractall(temp_path)
            
            # 保存SVG文件
            svg_filename = 'image1.svg'
            media_dir = temp_path / 'ppt' / 'media'
            media_dir.mkdir(parents=True, exist_ok=True)
            svg_path = media_dir / svg_filename
            
            with open(svg_path, 'w', encoding='utf-8') as f:
                f.write(svg_content)
            
            # 修改slide1.xml，添加SVG引用
            slide_path = temp_path / 'ppt' / 'slides' / 'slide1.xml'
            
            # 使用lxml解析XML
            parser = etree.XMLParser(remove_blank_text=True)
            tree = etree.parse(str(slide_path), parser)
            root = tree.getroot()
            
            # 定义命名空间
            namespaces = {
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
                'pic': 'http://schemas.openxmlformats.org/drawingml/2006/picture'
            }
            
            for prefix, uri in namespaces.items():
                etree.register_namespace(prefix, uri)
            
            # 获取spTree元素
            sp_tree = root.find('.//p:cSld/p:spTree', namespaces)
            
            if sp_tree is not None:
                # 转换单位为EMU (English Metric Units)
                from pptx.util import Emu
                x_emu = int(left)
                y_emu = int(top)
                cx_emu = int(width)
                cy_emu = int(height)
                
                # 创建图片元素
                pic_elem = self._create_svg_picture_element(
                    namespaces, x_emu, y_emu, cx_emu, cy_emu, 'rId2'
                )
                
                sp_tree.append(pic_elem)
            
            # 保存修改后的slide1.xml
            tree.write(str(slide_path), encoding='utf-8', xml_declaration=True, pretty_print=True)
            
            # 修改slide1.xml.rels，添加关系
            rels_path = temp_path / 'ppt' / 'slides' / '_rels' / 'slide1.xml.rels'
            rels_path.parent.mkdir(parents=True, exist_ok=True)
            
            if rels_path.exists():
                rels_tree = etree.parse(str(rels_path), parser)
                rels_root = rels_tree.getroot()
            else:
                # 创建新的rels文件
                rels_root = etree.Element(
                    '{http://schemas.openxmlformats.org/package/2006/relationships}Relationships'
                )
            
            # 添加SVG关系
            rel_elem = etree.SubElement(
                rels_root,
                '{http://schemas.openxmlformats.org/package/2006/relationships}Relationship',
                Id='rId2',
                Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/image',
                Target=f'../media/{svg_filename}'
            )
            
            # 保存rels文件
            rels_tree = etree.ElementTree(rels_root)
            rels_tree.write(str(rels_path), encoding='utf-8', xml_declaration=True, pretty_print=True)
            
            # 修改[Content_Types].xml，添加SVG类型
            content_types_path = temp_path / '[Content_Types].xml'
            ct_tree = etree.parse(str(content_types_path), parser)
            ct_root = ct_tree.getroot()
            
            # 检查是否已经有svg类型
            svg_exists = ct_root.find(
                './/{http://schemas.openxmlformats.org/package/2006/content-types}Default[@Extension="svg"]'
            )
            
            if svg_exists is None:
                # 添加SVG类型
                default_elem = etree.SubElement(
                    ct_root,
                    '{http://schemas.openxmlformats.org/package/2006/content-types}Default',
                    Extension='svg',
                    ContentType='image/svg+xml'
                )
            
            ct_tree.write(str(content_types_path), encoding='utf-8', xml_declaration=True, pretty_print=True)
            
            # 重新打包为PPTX
            with zipfile.ZipFile(output_pptx, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for file_path in temp_path.rglob('*'):
                    if file_path.is_file():
                        arcname = file_path.relative_to(temp_path)
                        zipf.write(file_path, arcname)
    
    def _create_svg_picture_element(self, namespaces, x, y, cx, cy, rid):
        """
        创建SVG图片元素的XML结构
        """
        from lxml import etree
        
        # 创建<p:pic>元素
        pic = etree.Element(f"{{{namespaces['p']}}}pic")
        
        # <p:nvPicPr>
        nvPicPr = etree.SubElement(pic, f"{{{namespaces['p']}}}nvPicPr")
        cNvPr = etree.SubElement(nvPicPr, f"{{{namespaces['p']}}}cNvPr", id="2", name="SVG Image")
        cNvPicPr = etree.SubElement(nvPicPr, f"{{{namespaces['p']}}}cNvPicPr")
        nvPr = etree.SubElement(nvPicPr, f"{{{namespaces['p']}}}nvPr")
        
        # <p:blipFill>
        blipFill = etree.SubElement(pic, f"{{{namespaces['p']}}}blipFill")
        blip = etree.SubElement(
            blipFill,
            f"{{{namespaces['a']}}}blip",
            {f"{{{namespaces['r']}}}embed": rid}
        )
        stretch = etree.SubElement(blipFill, f"{{{namespaces['a']}}}stretch")
        fillRect = etree.SubElement(stretch, f"{{{namespaces['a']}}}fillRect")
        
        # <p:spPr>
        spPr = etree.SubElement(pic, f"{{{namespaces['p']}}}spPr")
        xfrm = etree.SubElement(spPr, f"{{{namespaces['a']}}}xfrm")
        off = etree.SubElement(xfrm, f"{{{namespaces['a']}}}off", x=str(x), y=str(y))
        ext = etree.SubElement(xfrm, f"{{{namespaces['a']}}}ext", cx=str(cx), cy=str(cy))
        
        prstGeom = etree.SubElement(spPr, f"{{{namespaces['a']}}}prstGeom", prst="rect")
        avLst = etree.SubElement(prstGeom, f"{{{namespaces['a']}}}avLst")
        
        return pic
    
    def cleanup(self, filepath: str):
        """
        清理临时文件
        
        Args:
            filepath: 要删除的文件路径
        """
        try:
            if os.path.exists(filepath):
                os.remove(filepath)
        except Exception as e:
            print(f"清理文件失败: {e}")


# 单例实例
_export_service = None


def get_export_service() -> ExportService:
    """获取导出服务实例"""
    global _export_service
    if _export_service is None:
        _export_service = ExportService()
    return _export_service
