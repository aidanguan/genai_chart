#!/usr/bin/env python3
"""检查PPTX文件内容"""

from pptx import Presentation
import sys

pptx_path = r"c:\AI\genai_chart-1\backend\temp\exports\infographic_1764236291567.pptx"

try:
    prs = Presentation(pptx_path)
    
    print(f"✅ PPTX文件打开成功")
    print(f"📊 幻灯片数量: {len(prs.slides)}")
    print(f"📐 幻灯片尺寸: {prs.slide_width/914400:.2f}\" x {prs.slide_height/914400:.2f}\" (英寸)")
    print(f"📐 比例: {(prs.slide_width/prs.slide_height):.2f}:1")
    print()
    
    for i, slide in enumerate(prs.slides):
        print(f"\n{'='*60}")
        print(f"幻灯片 #{i+1}")
        print(f"{'='*60}")
        print(f"形状数量: {len(slide.shapes)}")
        
        for j, shape in enumerate(slide.shapes):
            print(f"\n  形状 #{j+1}: {shape.shape_type}")
            
            # 检查是否为图片
            if hasattr(shape, 'image'):
                print(f"    类型: 图片")
                print(f"    位置: ({shape.left/914400:.2f}\", {shape.top/914400:.2f}\")")
                print(f"    尺寸: {shape.width/914400:.2f}\" x {shape.height/914400:.2f}\"")
            
            # 检查是否为文本框
            if hasattr(shape, 'text_frame'):
                print(f"    类型: 文本框")
                print(f"    位置: ({shape.left/914400:.2f}\", {shape.top/914400:.2f}\")")
                print(f"    尺寸: {shape.width/914400:.2f}\" x {shape.height/914400:.2f}\"")
                print(f"    文本: \"{shape.text}\"")
                
                if shape.text:
                    for para in shape.text_frame.paragraphs:
                        if para.runs:
                            run = para.runs[0]
                            print(f"    字体大小: {run.font.size}")
                            print(f"    字体大小(pt): {run.font.size.pt if run.font.size else 'N/A'} pt")
                            print(f"    字体颜色: {run.font.color}")
                            if hasattr(run.font.color, 'rgb'):
                                print(f"    RGB: {run.font.color.rgb}")
                            print(f"    字体名称: {run.font.name if run.font.name else 'N/A'}")
                            print(f"    是否粗体: {run.font.bold}")
                            break
    
    print(f"\n{'='*60}")
    print(f"✅ 检查完成")
    
except Exception as e:
    print(f"❌ 错误: {e}")
    import traceback
    traceback.print_exc()
    sys.exit(1)
